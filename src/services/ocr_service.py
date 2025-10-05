# -*- coding: utf-8 -*-
"""
OCR 服务
使用 Doubao Vision 识别文档中的文字
支持 PDF, PPT, DOCX, DOC 等格式
"""

import logging
import os
import tempfile
from pathlib import Path
from typing import List, Dict, Any, Optional
from io import BytesIO

# 文档处理库
try:
    from pdf2image import convert_from_path, convert_from_bytes
    PDF_AVAILABLE = True
except ImportError:
    PDF_AVAILABLE = False
    logging.warning("pdf2image 未安装，PDF 转换功能不可用")

try:
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False
    logging.warning("python-pptx 未安装，PPT 处理功能不可用")

try:
    from docx import Document
    from docx.oxml.table import CT_Tbl
    from docx.oxml.text.paragraph import CT_P
    from docx.table import _Cell, Table
    from docx.text.paragraph import Paragraph
    DOCX_AVAILABLE = True
except ImportError:
    DOCX_AVAILABLE = False
    logging.warning("python-docx 未安装，DOCX 处理功能不可用")

from PIL import Image
import base64

from ..llms.providers.doubao_llm import DoubaoProvider
from ..config.config_loader import get_settings

logger = logging.getLogger(__name__)
settings = get_settings()


class OCRService:
    """OCR 服务类"""
    
    def __init__(self):
        """初始化 OCR 服务"""
        self.doubao_provider = None
        self._init_doubao()
    
    def _init_doubao(self):
        """初始化 Doubao Provider"""
        try:
            self.doubao_provider = DoubaoProvider(
                model_name=settings.doubao_vision_model,
                api_key=settings.doubao_api_key,
                base_url=settings.doubao_base_url
            )
            logger.info("Doubao Vision OCR 服务已初始化")
        except Exception as e:
            logger.error(f"初始化 Doubao Vision 失败: {e}")
    
    async def recognize_image(self, image_path: str) -> Dict[str, Any]:
        """
        识别单张图片中的文字
        
        Args:
            image_path: 图片路径
        
        Returns:
            识别结果
        """
        try:
            # 读取图片并转换为 base64
            with open(image_path, 'rb') as f:
                image_data = f.read()
            
            image_base64 = base64.b64encode(image_data).decode('utf-8')
            
            # 构建消息
            messages = [
                {
                    "role": "user",
                    "content": [
                        {
                            "type": "image_url",
                            "image_url": {
                                "url": f"data:image/jpeg;base64,{image_base64}"
                            }
                        },
                        {
                            "type": "text",
                            "text": "请识别图片中的所有文字内容，保持原有格式和结构。"
                        }
                    ]
                }
            ]
            
            # 调用 Doubao Vision
            response = await self.doubao_provider.generate(messages)
            
            return {
                "success": True,
                "text": response.content,
                "image_path": image_path
            }
        
        except Exception as e:
            logger.error(f"图片识别失败: {e}")
            return {
                "success": False,
                "error": str(e),
                "image_path": image_path
            }
    
    async def pdf_to_text(self, pdf_path: str) -> Dict[str, Any]:
        """
        将 PDF 转换为文本
        
        Args:
            pdf_path: PDF 文件路径
        
        Returns:
            提取的文本内容
        """
        if not PDF_AVAILABLE:
            return {
                "success": False,
                "error": "pdf2image 未安装"
            }
        
        try:
            logger.info(f"开始处理 PDF: {pdf_path}")
            
            # 转换 PDF 为图片
            images = convert_from_path(pdf_path, dpi=200)
            
            logger.info(f"PDF 共 {len(images)} 页")
            
            # 识别每一页
            pages_text = []
            for i, image in enumerate(images, 1):
                logger.info(f"识别第 {i}/{len(images)} 页")
                
                # 保存临时图片
                temp_image_path = tempfile.mktemp(suffix='.jpg')
                image.save(temp_image_path, 'JPEG')
                
                # 识别文字
                result = await self.recognize_image(temp_image_path)
                
                if result["success"]:
                    pages_text.append({
                        "page": i,
                        "text": result["text"]
                    })
                
                # 删除临时文件
                os.remove(temp_image_path)
            
            # 合并所有页面的文本
            full_text = "\n\n".join([
                f"=== 第 {page['page']} 页 ===\n{page['text']}"
                for page in pages_text
            ])
            
            return {
                "success": True,
                "text": full_text,
                "pages": pages_text,
                "total_pages": len(images)
            }
        
        except Exception as e:
            logger.error(f"PDF 处理失败: {e}")
            return {
                "success": False,
                "error": str(e)
            }
    
    async def ppt_to_text(self, ppt_path: str) -> Dict[str, Any]:
        """
        将 PPT 转换为文本
        
        Args:
            ppt_path: PPT 文件路径
        
        Returns:
            提取的文本内容
        """
        if not PPTX_AVAILABLE:
            return {
                "success": False,
                "error": "python-pptx 未安装"
            }
        
        try:
            logger.info(f"开始处理 PPT: {ppt_path}")
            
            prs = Presentation(ppt_path)
            
            logger.info(f"PPT 共 {len(prs.slides)} 页")
            
            slides_text = []
            
            for i, slide in enumerate(prs.slides, 1):
                logger.info(f"处理第 {i}/{len(prs.slides)} 页")
                
                # 提取文本内容
                text_content = []
                
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text_content.append(shape.text)
                
                # 如果有图片，截图整页进行 OCR
                if any(shape.shape_type == MSO_SHAPE_TYPE.PICTURE for shape in slide.shapes):
                    # 导出幻灯片为图片
                    temp_image_path = tempfile.mktemp(suffix='.png')
                    
                    # 注意：python-pptx 不直接支持导出图片
                    # 这里需要使用其他方法，如 LibreOffice 或 unoconv
                    # 暂时跳过图片 OCR
                    logger.warning(f"第 {i} 页包含图片，但暂不支持图片 OCR")
                
                slide_text = "\n".join(text_content)
                
                slides_text.append({
                    "slide": i,
                    "text": slide_text
                })
            
            # 合并所有幻灯片的文本
            full_text = "\n\n".join([
                f"=== 第 {slide['slide']} 页 ===\n{slide['text']}"
                for slide in slides_text
            ])
            
            return {
                "success": True,
                "text": full_text,
                "slides": slides_text,
                "total_slides": len(prs.slides)
            }
        
        except Exception as e:
            logger.error(f"PPT 处理失败: {e}")
            return {
                "success": False,
                "error": str(e)
            }
    
    async def docx_to_text(self, docx_path: str) -> Dict[str, Any]:
        """
        将 DOCX 转换为文本
        
        Args:
            docx_path: DOCX 文件路径
        
        Returns:
            提取的文本内容
        """
        if not DOCX_AVAILABLE:
            return {
                "success": False,
                "error": "python-docx 未安装"
            }
        
        try:
            logger.info(f"开始处理 DOCX: {docx_path}")
            
            doc = Document(docx_path)
            
            # 提取所有段落和表格
            full_text = []
            
            for element in doc.element.body:
                if isinstance(element, CT_P):
                    # 段落
                    paragraph = Paragraph(element, doc)
                    if paragraph.text.strip():
                        full_text.append(paragraph.text)
                
                elif isinstance(element, CT_Tbl):
                    # 表格
                    table = Table(element, doc)
                    table_text = []
                    for row in table.rows:
                        row_text = " | ".join([cell.text for cell in row.cells])
                        table_text.append(row_text)
                    full_text.append("\n".join(table_text))
            
            text = "\n\n".join(full_text)
            
            return {
                "success": True,
                "text": text,
                "paragraphs": len(doc.paragraphs),
                "tables": len(doc.tables)
            }
        
        except Exception as e:
            logger.error(f"DOCX 处理失败: {e}")
            return {
                "success": False,
                "error": str(e)
            }
    
    async def process_document(self, file_path: str) -> Dict[str, Any]:
        """
        处理文档（自动识别格式）
        
        Args:
            file_path: 文件路径
        
        Returns:
            提取的文本内容
        """
        file_ext = Path(file_path).suffix.lower()
        
        if file_ext == '.pdf':
            return await self.pdf_to_text(file_path)
        elif file_ext in ['.ppt', '.pptx']:
            return await self.ppt_to_text(file_path)
        elif file_ext in ['.doc', '.docx']:
            return await self.docx_to_text(file_path)
        elif file_ext in ['.jpg', '.jpeg', '.png', '.bmp']:
            result = await self.recognize_image(file_path)
            return {
                "success": result["success"],
                "text": result.get("text", ""),
                "error": result.get("error")
            }
        else:
            return {
                "success": False,
                "error": f"不支持的文件格式: {file_ext}"
            }


# 全局实例
_ocr_service = None


def get_ocr_service() -> OCRService:
    """获取 OCR 服务实例"""
    global _ocr_service
    if _ocr_service is None:
        _ocr_service = OCRService()
    return _ocr_service
