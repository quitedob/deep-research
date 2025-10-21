from io import BytesIO
from typing import Dict, Any, List, Optional

try:
    from pptx import Presentation
    from pptx.util import Pt
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False
    Presentation = None
    Pt = None

from .base import Exporter


class PPTXExporter(Exporter):
    """生成真实 .pptx 幻灯片文件。"""

    def __init__(self, template: Optional[str] = None):
        if not PPTX_AVAILABLE:
            raise ImportError("python-pptx 库未安装，无法使用PPTXExporter")
        self.template = template

    def _parse_content_to_slides(self, content: str, title: str) -> List[Dict[str, Any]]:
        lines = [ln.rstrip() for ln in content.splitlines()]
        slides: List[Dict[str, Any]] = []

        current_title = title or "演示文稿"
        current_body: List[str] = []

        for ln in lines:
            if ln.startswith("#"):
                if current_body:
                    slides.append({
                        "type": "content",
                        "title": current_title,
                        "content": "\n".join(current_body).strip(),
                    })
                    current_body = []
                current_title = ln.lstrip("#").strip() or current_title
            else:
                current_body.append(ln)

        if current_body:
            slides.append({
                "type": "content",
                "title": current_title,
                "content": "\n".join(current_body).strip(),
            })

        if not slides:
            slides.append({"type": "content", "title": title or "演示文稿", "content": ""})
        return slides

    def export(self, content: str, title: str, **kwargs) -> bytes:
        if self.template:
            try:
                prs = Presentation(self.template)
            except Exception:
                prs = Presentation()
        else:
            prs = Presentation()

        # 封面
        title_layout = prs.slide_layouts[0]
        cover = prs.slides.add_slide(title_layout)
        cover.shapes.title.text = title or "演示文稿"
        if len(cover.placeholders) > 1:
            cover.placeholders[1].text = kwargs.get("subtitle", "由 Deep Research 自动生成")
            cover.placeholders[1].text_frame.paragraphs[0].font.size = Pt(18)

        # 内容
        for s in self._parse_content_to_slides(content, title):
            layout = prs.slide_layouts[1]
            slide = prs.slides.add_slide(layout)
            slide.shapes.title.text = s.get("title") or "内容"
            title_par = slide.shapes.title.text_frame.paragraphs[0]
            title_par.font.size = Pt(28)
            title_par.font.bold = True

            body_tf = slide.placeholders[1].text_frame
            body_tf.clear()
            body_text = s.get("content", "")
            paragraphs = [p for p in body_text.split("\n\n") if p.strip()]
            if not paragraphs:
                body_tf.add_paragraph().text = ""
            for para in paragraphs:
                p = body_tf.add_paragraph()
                p.text = para.strip()
                p.font.size = Pt(18)

        bio = BytesIO()
        prs.save(bio)
        return bio.getvalue()


