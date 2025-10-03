# -*- coding: utf-8 -*-
"""
PPTX渲染器

将DSL（XML/JSON）转换为PPTX文件。
"""

import logging
import xml.etree.ElementTree as ET
from pathlib import Path
from typing import Dict, List, Optional, Any
from io import BytesIO

try:
    from pptx import Presentation
    from pptx.util import Pt, Inches
    from pptx.enum.text import PP_ALIGN
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False
    logger.warning("python-pptx未安装，PPT渲染功能不可用")

logger = logging.getLogger(__name__)


class PPTXRenderer:
    """PPTX渲染器类"""

    def __init__(self):
        """初始化渲染器"""
        if not PPTX_AVAILABLE:
            raise ImportError("python-pptx库未安装，无法使用PPTXRenderer")

        self.supported_layouts = [
            "TITLE", "BULLETS", "COLUMNS", "IMAGE",
            "ICONS", "TIMELINE", "CHART"
        ]

    async def render_dsl_to_pptx(
        self,
        dsl_content: str,
        title: str,
        template: Optional[str] = None,
        output_dir: Path = None,
        presentation_id: str = None
    ) -> Path:
        """
        将DSL渲染为PPTX文件

        参数:
            dsl_content: DSL内容（XML格式）
            title: 演示文稿标题
            template: 模板名称（可选）
            output_dir: 输出目录
            presentation_id: 演示文稿ID

        返回:
            生成的PPTX文件路径
        """
        try:
            # 1. 解析DSL
            slides_data = self._parse_dsl(dsl_content)

            # 2. 创建演示文稿
            if template:
                prs = self._load_template(template)
            else:
                prs = Presentation()

            # 3. 渲染每一页
            for slide_data in slides_data:
                self._render_slide(prs, slide_data)

            # 4. 保存文件
            if output_dir is None:
                output_dir = Path("./output_reports/ppt")
            output_dir.mkdir(parents=True, exist_ok=True)

            filename = f"{presentation_id or 'presentation'}_{title[:30]}.pptx"
            # 清理文件名中的非法字符
            filename = "".join(c for c in filename if c.isalnum() or c in "._- ")
            output_path = output_dir / filename

            prs.save(str(output_path))
            logger.info(f"PPTX文件已保存: {output_path}")

            return output_path

        except Exception as e:
            logger.error(f"渲染PPTX失败: {str(e)}", exc_info=True)
            raise

    def _parse_dsl(self, dsl_content: str) -> List[Dict[str, Any]]:
        """
        解析DSL内容

        参数:
            dsl_content: DSL内容

        返回:
            幻灯片数据列表
        """
        slides = []

        try:
            # 尝试解析XML
            # 清理可能的markdown代码块标记
            dsl_content = dsl_content.strip()
            if dsl_content.startswith("```xml"):
                dsl_content = dsl_content[6:]
            if dsl_content.startswith("```"):
                dsl_content = dsl_content[3:]
            if dsl_content.endswith("```"):
                dsl_content = dsl_content[:-3]
            dsl_content = dsl_content.strip()

            root = ET.fromstring(dsl_content)

            # 提取SECTION元素
            for section in root.findall(".//SECTION"):
                slide_data = self._parse_section(section)
                if slide_data:
                    slides.append(slide_data)

            if not slides:
                logger.warning("未找到SECTION元素，尝试备用解析方法")
                slides = self._parse_fallback(dsl_content)

        except ET.ParseError as e:
            logger.error(f"XML解析失败: {str(e)}，使用备用解析方法")
            slides = self._parse_fallback(dsl_content)

        return slides

    def _parse_section(self, section: ET.Element) -> Optional[Dict[str, Any]]:
        """解析单个SECTION元素"""
        layout = section.get("layout", "BULLETS")

        slide_data = {
            "layout": layout,
            "title": "",
            "content": []
        }

        # 提取标题
        title_elem = section.find("TITLE")
        if title_elem is not None and title_elem.text:
            slide_data["title"] = title_elem.text.strip()

        # 根据布局类型提取内容
        if layout == "TITLE":
            subtitle_elem = section.find("SUBTITLE")
            if subtitle_elem is not None and subtitle_elem.text:
                slide_data["subtitle"] = subtitle_elem.text.strip()

        elif layout == "BULLETS":
            content_elem = section.find("CONTENT")
            if content_elem is not None:
                bullets = content_elem.findall("BULLET")
                slide_data["content"] = [b.text.strip() for b in bullets if b.text]

        elif layout == "COLUMNS":
            content_elem = section.find("CONTENT")
            if content_elem is not None:
                columns = []
                for col in content_elem.findall("COLUMN"):
                    heading = col.find("HEADING")
                    text = col.find("TEXT")
                    columns.append({
                        "heading": heading.text.strip() if heading is not None and heading.text else "",
                        "text": text.text.strip() if text is not None and text.text else ""
                    })
                slide_data["columns"] = columns

        # 提取图像查询（如果有）
        image_query = section.find("IMAGE_QUERY")
        if image_query is not None and image_query.text:
            slide_data["image_query"] = image_query.text.strip()

        return slide_data

    def _parse_fallback(self, content: str) -> List[Dict[str, Any]]:
        """备用解析方法：从纯文本中提取内容"""
        slides = []

        # 按行分割
        lines = content.split("\n")
        current_slide = None

        for line in lines:
            line = line.strip()
            if not line:
                continue

            # 检测标题（以#开头）
            if line.startswith("#"):
                if current_slide:
                    slides.append(current_slide)

                current_slide = {
                    "layout": "BULLETS",
                    "title": line.lstrip("#").strip(),
                    "content": []
                }

            # 检测列表项
            elif line.startswith("-") or line.startswith("*") or line.startswith("•"):
                if current_slide:
                    current_slide["content"].append(line.lstrip("-*• ").strip())

        # 添加最后一个slide
        if current_slide:
            slides.append(current_slide)

        # 如果没有解析到任何内容，创建一个默认slide
        if not slides:
            slides.append({
                "layout": "TITLE",
                "title": "演示文稿",
                "subtitle": "自动生成"
            })

        return slides

    def _render_slide(self, prs: Presentation, slide_data: Dict[str, Any]):
        """
        渲染单个幻灯片

        参数:
            prs: Presentation对象
            slide_data: 幻灯片数据
        """
        layout_type = slide_data.get("layout", "BULLETS")

        if layout_type == "TITLE":
            self._render_title_slide(prs, slide_data)
        elif layout_type == "BULLETS":
            self._render_bullets_slide(prs, slide_data)
        elif layout_type == "COLUMNS":
            self._render_columns_slide(prs, slide_data)
        else:
            # 其他布局暂时使用BULLETS布局
            self._render_bullets_slide(prs, slide_data)

    def _render_title_slide(self, prs: Presentation, slide_data: Dict[str, Any]):
        """渲染标题页"""
        slide_layout = prs.slide_layouts[0]  # 标题布局
        slide = prs.slides.add_slide(slide_layout)

        title = slide.shapes.title
        title.text = slide_data.get("title", "")

        if len(slide.placeholders) > 1:
            subtitle = slide.placeholders[1]
            subtitle.text = slide_data.get("subtitle", "")

    def _render_bullets_slide(self, prs: Presentation, slide_data: Dict[str, Any]):
        """渲染项目符号页"""
        slide_layout = prs.slide_layouts[1]  # 标题和内容布局
        slide = prs.slides.add_slide(slide_layout)

        title = slide.shapes.title
        title.text = slide_data.get("title", "")

        body_shape = slide.placeholders[1]
        text_frame = body_shape.text_frame
        text_frame.clear()

        content = slide_data.get("content", [])
        if not content:
            content = ["内容待补充"]

        for item in content:
            p = text_frame.add_paragraph()
            p.text = item
            p.level = 0
            p.font.size = Pt(18)

    def _render_columns_slide(self, prs: Presentation, slide_data: Dict[str, Any]):
        """渲染多列布局页"""
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)

        title = slide.shapes.title
        title.text = slide_data.get("title", "")

        # 简化处理：将列内容作为项目符号显示
        body_shape = slide.placeholders[1]
        text_frame = body_shape.text_frame
        text_frame.clear()

        columns = slide_data.get("columns", [])
        for col in columns:
            heading = col.get("heading", "")
            text = col.get("text", "")

            if heading:
                p = text_frame.add_paragraph()
                p.text = heading
                p.level = 0
                p.font.bold = True
                p.font.size = Pt(20)

            if text:
                p = text_frame.add_paragraph()
                p.text = text
                p.level = 1
                p.font.size = Pt(16)

    def _load_template(self, template_name: str) -> Presentation:
        """加载模板"""
        # 这里可以根据template_name加载不同的模板
        # 暂时返回默认模板
        return Presentation()


# 全局实例
_renderer: Optional[PPTXRenderer] = None


def get_renderer() -> PPTXRenderer:
    """获取渲染器实例（单例模式）"""
    global _renderer
    if _renderer is None:
        _renderer = PPTXRenderer()
    return _renderer
